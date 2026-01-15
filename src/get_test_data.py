import os
import random
import xlsxwriter
from docx import Document
from pptx import Presentation
from pptx.util import Inches

# CONFIG
BENIGN_DEST = "../data/test/benign"
LIMIT = 50  # How many files to make

def generate_benign_suite(limit=50):
    if not os.path.exists(BENIGN_DEST):
        os.makedirs(BENIGN_DEST)

    print(f"[*] Generating {limit} benign files (Word, Excel, PowerPoint)...")

    # The list of formats we want to mimic
    formats = ['docx', 'docm', 'xlsx', 'xlsm', 'pptx', 'pptm']

    for i in range(limit):
        # Pick a random format
        ftype = random.choice(formats)
        fname = f"safe_sample_{i}.{ftype}"
        fpath = os.path.join(BENIGN_DEST, fname)
        
        try:
            # --- WORD (.docx, .docm) ---
            if ftype in ['docx', 'docm']:
                doc = Document()
                doc.add_heading(f'Benign Report {i}', 0)
                doc.add_paragraph("This file contains standard XML structure.")
                # Add a table to make the internal XML complex/realistic
                table = doc.add_table(rows=2, cols=2)
                table.cell(0, 0).text = "Safe"
                table.cell(1, 1).text = "Data"
                doc.save(fpath)

            # --- EXCEL (.xlsx, .xlsm) ---
            elif ftype in ['xlsx', 'xlsm']:
                # XlsxWriter supports saving as xlsm even without macros
                workbook = xlsxwriter.Workbook(fpath)
                worksheet = workbook.add_worksheet()
                worksheet.write('A1', 'Financial Safe Data')
                worksheet.write('B1', random.randint(100, 9999))
                # Add a chart to add structural complexity
                chart = workbook.add_chart({'type': 'column'})
                worksheet.insert_chart('C1', chart)
                workbook.close()

            # --- POWERPOINT (.pptx, .pptm) ---
            elif ftype in ['pptx', 'pptm']:
                prs = Presentation()
                # Add a title slide
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = f"Safe Presentation {i}"
                subtitle.text = "Generated for AI Training"
                
                # Add a content slide with bullet points
                bullet_slide_layout = prs.slide_layouts[1]
                slide2 = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide2.shapes
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                title_shape.text = "Agenda"
                tf = body_shape.text_frame
                tf.text = "Item 1: No Macros"
                p = tf.add_paragraph()
                p.text = "Item 2: Safe Structure"
                
                prs.save(fpath)

            print(f"    [+] Created: {fname}")

        except Exception as e:
            print(f"    [-] Failed to create {fname}: {e}")

    print(f"[*] Done. Files saved to {BENIGN_DEST}")

if __name__ == "__main__":
    generate_benign_suite(LIMIT)